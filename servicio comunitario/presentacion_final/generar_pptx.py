import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# Definir la paleta de colores premium
AZUL_OSCURO = RGBColor(0, 51, 102)
AZUL_CLARO = RGBColor(0, 102, 204)
BLANCO = RGBColor(255, 255, 255)
GRIS_SUAVE = RGBColor(245, 245, 245)
GRIS_TEXTO = RGBColor(64, 64, 64)

def apply_slide_template(slide, is_dark=False):
    if is_dark:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = AZUL_OSCURO # Corregido: fore_color en lugar de foreground_color
    else:
        # Agregar una barra lateral decorativa
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = AZUL_OSCURO # Corregido
        shape.line.fill.background() # Sin borde

def add_styled_slide(title, points, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Usar layout en blanco para control total
    apply_slide_template(slide)
    
    # Titulo
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8.5)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = AZUL_OSCURO
    
    # Contenido (ajustar si hay imagen)
    content_width = Inches(5) if image_path else Inches(8)
    txContent = slide.shapes.add_textbox(Inches(1), Inches(1.8), content_width, Inches(4.5))
    tf_c = txContent.text_frame
    tf_c.word_wrap = True
    
    for i, point in enumerate(points):
        p_c = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        p_c.text = f"• {point}"
        p_c.font.size = Pt(18)
        p_c.font.color.rgb = GRIS_TEXTO
        p_c.space_after = Pt(12)

    # Imagen con posicionamiento estético
    if image_path and os.path.exists(image_path):
        img_left = Inches(6.5)
        img_top = Inches(1.8)
        img_width = Inches(3)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

# --- Directorio de Imágenes ---
img_dir = "/home/sergio/Documentos/servicio comunitario/"
images = [f for f in os.listdir(img_dir) if f.endswith('.jpeg') or f.endswith('.jpg')]
images.sort()

# --- PORTADA PREMIUM ---
slide_p = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_template(slide_p, is_dark=True)

title_p = slide_p.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
tf_p = title_p.text_frame
p_p = tf_p.paragraphs[0]
p_p.text = "INFORME FINAL DE SERVICIO COMUNITARIO"
p_p.alignment = PP_ALIGN.CENTER
p_p.font.bold = True
p_p.font.size = Pt(44)
p_p.font.color.rgb = BLANCO

subtitle_p = slide_p.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf_s = subtitle_p.text_frame
p_s = tf_s.paragraphs[0]
p_s.text = "Sergio Malavé | IUJO - Informática | 2026"
p_s.alignment = PP_ALIGN.CENTER
p_s.font.size = Pt(20)
p_s.font.color.rgb = BLANCO

# --- DIAPOSITIVAS DE CONTENIDO ---
content_slides = [
    ("Introducción", ["Sistematización del proyecto de Responsabilidad Social.", "Alianza estratégica con SUPERATEC Propatria.", "Objetivo: Reducción de la brecha digital en la comunidad.", "Duración: 3 meses de impacto directo."]),
    ("La Comunidad", ["Ubicación: Sede SUPERATEC, Centro Comunal Catia Cecca.", "Población: 15 participantes activos (jóvenes y adultos).", "Composición: 10 hombres y 5 mujeres.", "Entorno: Comunidad de Propatria, Caracas."]),
    ("Problemática", ["Desconexión entre acceso básico y formación avanzada.", "Barreras psicológicas y baja autoeficacia tecnológica.", "Estereotipos que limitan el aprendizaje en mujeres y adultos mayores.", "Necesidad de rutas de aprendizaje claras y actualizadas."]),
    ("Estrategia Metodológica", ["Modelo: Aprendizaje por Proyectos (ABP).", "Formato: Híbrido (Sesiones Teórico-Prácticas).", "Estructura: 4 horas semanales de inmersión técnica.", "Evaluación: Cualitativa y cuantitativa por desempeño."]),
    ("Actividades Técnicas", ["Fundamentos de la Web y Estructura HTML5.", "Maquetación y Estilos con CSS3 (Responsive Design).", "Lógica de Programación y Variables en JavaScript.", "Desarrollo de proyectos finales funcionales."]),
    ("Resultados Clave", ["92% de asistencia promedio.", "15 sitios web personales creados desde cero.", "85% de los participantes aprobados con excelencia.", "95% nivel de satisfacción en la comunidad."]),
    ("Perfil Académico", ["Aplicación de Programación I, II y Diseño Web.", "Sistemas Operativos e Ingeniería de Software.", "Ética Profesional y Compromiso Social IUJO.", "Consolidación de la vocación docente-tecnológica."]),
    ("Conclusiones", ["Éxito en la reducción de barreras tecnológicas.", "Importancia de las alianzas universidad-comunidad.", "Necesidad de continuidad en la formación digital.", "Impacto positivo en la empleabilidad local."])
]

for i, (title, points) in enumerate(content_slides):
    img = os.path.join(img_dir, images[i % len(images)]) if images else None
    add_styled_slide(title, points, image_path=img)

# Galería Final
for img_file in images:
    slide_g = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_template(slide_g)
    slide_g.shapes.add_picture(os.path.join(img_dir, img_file), Inches(1.5), Inches(1), height=Inches(5.5))
    
# Cierre
slide_c = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_template(slide_c, is_dark=True)
txt_c = slide_c.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
p_c = txt_c.text_frame.paragraphs[0]
p_c.text = "¡MUCHAS GRACIAS POR SU ATENCIÓN!"
p_c.alignment = PP_ALIGN.CENTER
p_c.font.bold = True
p_c.font.size = Pt(40)
p_c.font.color.rgb = BLANCO

# Guardar
output_path = "/home/sergio/Documentos/servicio comunitario/presentacion_final/presentacion_final_sergio.pptx"
prs.save(output_path)
print(f"Versión PREMIUM generada exitosamente en {output_path}")
